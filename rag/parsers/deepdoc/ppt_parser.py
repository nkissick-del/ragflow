#
#  Copyright 2025 The InfiniFlow Authors. All Rights Reserved.
#
#  Licensed under the Apache License, Version 2.0 (the "License");
#  you may not use this file except in compliance with the License.
#  You may obtain a copy of the License at
#
#      http://www.apache.org/licenses/LICENSE-2.0
#
#  Unless required by applicable law or agreed to in writing, software
#  distributed under the License is distributed on an "AS IS" BASIS,
#  WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
#  See the License for the specific language governing permissions and
#  limitations under the License.
#

import logging
from io import BytesIO
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


class RAGFlowPptParser:
    def __init__(self):
        pass

    def __get_bulleted_text(self, paragraph):
        is_bulleted = bool(paragraph._p.xpath("./a:pPr/a:buChar")) or bool(paragraph._p.xpath("./a:pPr/a:buAutoNum")) or bool(paragraph._p.xpath("./a:pPr/a:buBlip"))
        if is_bulleted:
            return f"{'  ' * paragraph.level}• {paragraph.text}"
        else:
            return paragraph.text

    def __extract(self, shape):
        try:
            # First try to get text content
            if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                text_frame = shape.text_frame
                texts = []
                for paragraph in text_frame.paragraphs:
                    if paragraph.text.strip():
                        texts.append(self.__get_bulleted_text(paragraph))
                return "\n".join(texts)

            # Safely get shape_type
            try:
                shape_type = shape.shape_type
            except NotImplementedError:
                # If shape_type is not available, try to get text content
                if hasattr(shape, "text"):
                    return shape.text.strip()
                return ""

            # Handle table
            if shape_type == MSO_SHAPE_TYPE.TABLE:
                tb = shape.table
                rows = []
                for i in range(1, len(tb.rows)):
                    row_cells = []
                    for j in range(len(tb.columns)):
                        cell = tb.cell(i, j)
                        if cell.text.strip():
                            row_cells.append(tb.cell(0, j).text + ": " + cell.text)
                    if row_cells:
                        rows.append("; ".join(row_cells))
                return "\n".join(rows)

            # Handle group shape
            if shape_type == MSO_SHAPE_TYPE.GROUP:
                texts = []
                for p in sorted(shape.shapes, key=lambda x: (x.top // 10, x.left)):
                    t = self.__extract(p)
                    if t:
                        texts.append(t)
                return "\n".join(texts)

            return ""

        except Exception as e:
            logging.error(f"Error processing shape: {str(e)}")
            return ""

    def __call__(self, fnm, from_page, to_page):
        ppt = Presentation(fnm) if isinstance(fnm, str) else Presentation(BytesIO(fnm))
        txts = []
        for i, slide in enumerate(ppt.slides):
            if i < from_page:
                continue
            if i >= to_page:
                break
            texts = []
            for shape in sorted(slide.shapes, key=lambda x: ((x.top if x.top is not None else 0) // 10, x.left if x.left is not None else 0)):
                txt = self.__extract(shape)
                if txt:
                    texts.append(txt)
            txts.append("\n".join(texts))

        return txts
